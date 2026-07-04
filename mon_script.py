# -*- coding: utf-8 -*-
import os
import threading
import tkinter as tk
from tkinter import filedialog, ttk, messagebox
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from datetime import date

# ---------- Couleurs du thème ----------
BG_COLOR = "#F5F1EA"  # Ivoire chaud
BUTTON_COLOR = "#D95F5F"  # Corail
BUTTON_HOVER = "#C64F4F"  # Corail foncé
TEXT_COLOR = "#243447"  # Bleu-gris profond
ENTRY_BG = "#FFFDF9"  # Blanc cassé
PPT_BG = "#F6F1EA"
PPT_PANEL = "#FFFDF9"
PPT_ACCENT = "#D95F5F"
PPT_ACCENT_2 = "#2F4B7C"
PPT_TEXT = "#243447"

# ---------- Fonctions Utilitaires ----------
def get_unique_filename(base="Chorale", ext=".pptx"):
    today = date.today().strftime("%d-%m-%Y")
    filename = f"{base}_{today}{ext}"
    counter = 1
    unique_name = filename
    while os.path.exists(unique_name):
        unique_name = f"{base}_{today}_{counter}{ext}"
        counter += 1
    return unique_name

def split_block_recursive(block, max_lines=12):
    if len(block) <= max_lines:
        return [block]
    mid = len(block) // 2
    return split_block_recursive(block[:mid], max_lines) + split_block_recursive(block[mid:], max_lines)

def hex_to_rgb(color_code):
    return RGBColor.from_string(color_code.replace("#", ""))

def add_styled_slide(prs, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = hex_to_rgb(PPT_BG)

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.58))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(PPT_ACCENT)
    top_bar.line.fill.background()

    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(1.1), Inches(0.14), Inches(4.7))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = hex_to_rgb(PPT_ACCENT_2)
    left_bar.line.fill.background()

    decorative = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.55), Inches(5.9), Inches(1.55), Inches(1.55))
    decorative.fill.solid()
    decorative.fill.fore_color.rgb = hex_to_rgb(PPT_ACCENT_2)
    decorative.line.fill.background()

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(0.98), Inches(8.45), Inches(5.0))
    panel.fill.solid()
    panel.fill.fore_color.rgb = hex_to_rgb(PPT_PANEL)
    panel.line.color.rgb = hex_to_rgb(PPT_ACCENT)
    panel.line.width = Pt(1.5)

    text_box = slide.shapes.add_textbox(Inches(1.25), Inches(1.35), Inches(7.7), Inches(4.3))
    text_frame = text_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    font_size = 32
    if len(lines) > 4:
        font_size = 28
    if len(lines) > 7:
        font_size = 24

    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.CENTER
        if paragraph.runs:
            run = paragraph.runs[0]
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = hex_to_rgb(PPT_TEXT)
            run.font.name = "Arial"

    return slide

def generate_pptx_from_lines(lines, max_lines=12):
    prs = Presentation()
    block = []

    for line in lines:
        line = line.strip()
        if line == "" and block:
            small_blocks = split_block_recursive(block, max_lines)
            for b in small_blocks:
                add_styled_slide(prs, b)
            block = []
        else:
            block.append(line)

    if block:
        small_blocks = split_block_recursive(block, max_lines)
        for b in small_blocks:
            add_styled_slide(prs, b)

    filename = get_unique_filename()
    prs.save(filename)
    return filename

# ---------- Actions ----------
def choose_file():
    filepath = filedialog.askopenfilename(title="Sélectionner un fichier de paroles",
                                          filetypes=[("Fichiers texte", "*.txt"), ("Tous les fichiers", "*.*")])
    if filepath:
        with open(filepath, "r", encoding="utf-8") as f:
            lines = f.readlines()
        text_box.delete("1.0", tk.END)
        text_box.insert("1.0", "".join(lines))
        output_file(lines)

def generate_from_text():
    text = text_box.get("1.0", tk.END).strip()
    if not text:
        messagebox.showwarning("⚠️ Attention", "Veuillez saisir du texte ou sélectionner un fichier.")
        return
    lines = text.split("\n")
    output_file(lines)

def on_enter(e):
    e.widget['background'] = BUTTON_HOVER

def on_leave(e):
    e.widget['background'] = BUTTON_COLOR

def show_success_popup(pptx_path):
    popup = tk.Toplevel(root)
    popup.title("Succès")
    popup.geometry("400x200")
    popup.configure(bg=BG_COLOR)
    tk.Label(popup, text="PowerPoint créé !", font=("Arial", 14, "bold"), bg=BG_COLOR, fg=TEXT_COLOR).pack(pady=20)
    tk.Label(popup, text=f"Fichier : {os.path.basename(pptx_path)}", bg=BG_COLOR, fg=TEXT_COLOR).pack(pady=10)

    btn_frame = tk.Frame(popup, bg=BG_COLOR)
    btn_frame.pack(pady=10)

    tk.Button(btn_frame, text="Ouvrir", command=lambda: os.startfile(pptx_path),
              bg=BUTTON_COLOR, fg="white", relief="flat", padx=15, pady=5,
              font=("Arial", 10, "bold")).pack(side=tk.LEFT, padx=10)
    tk.Button(btn_frame, text="Fermer", command=popup.destroy,
              bg=BUTTON_COLOR, fg="white", relief="flat", padx=15, pady=5,
              font=("Arial", 10, "bold")).pack(side=tk.LEFT, padx=10)

def show_error_popup(message):
    messagebox.showerror("Erreur", message)

def output_file(lines):
    # Fenêtre de chargement
    loading_popup = tk.Toplevel(root)
    loading_popup.title("Chargement...")
    loading_popup.geometry("300x120")
    loading_popup.resizable(False, False)
    loading_popup.configure(bg=BG_COLOR)
    tk.Label(loading_popup, text="🎵 Création du PowerPoint...", font=("Arial", 12), bg=BG_COLOR, fg=TEXT_COLOR).pack(expand=True)
    progress = ttk.Progressbar(loading_popup, mode='indeterminate', length=200)
    progress.pack(pady=10)
    progress.start()
    loading_popup.update()

    result = {"path": None, "error": None}

    def worker():
        try:
            result["path"] = generate_pptx_from_lines(lines, max_lines=8)
        except Exception as exc:
            result["error"] = exc

    def poll_worker():
        if worker_thread.is_alive():
            root.after(100, poll_worker)
            return

        loading_popup.destroy()
        if result["error"] is not None:
            show_error_popup(f"Impossible de créer le PowerPoint : {result['error']}")
            return

        show_success_popup(result["path"])

    worker_thread = threading.Thread(target=worker, daemon=True)
    worker_thread.start()
    poll_worker()

# ---------- Interface ----------
root = tk.Tk()
root.title("🎤 Générateur PowerPoint pour Chorale")
root.geometry("1000x700")
root.minsize(800, 500)
root.configure(bg=BG_COLOR)

# Cadre principal
main_frame = tk.Frame(root, bg=BG_COLOR)
main_frame.pack(expand=True, fill=tk.BOTH, padx=30, pady=20)

# Titre
title_label = tk.Label(main_frame, text="Générateur de PowerPoint Chorale", 
                       font=("Arial", 18, "bold"), fg=TEXT_COLOR, bg=BG_COLOR)
title_label.pack(pady=(0, 20))

# Bouton de sélection de fichier
btn_file = tk.Button(main_frame, text="📂 Sélectionner un fichier", command=choose_file,
                     bg=BUTTON_COLOR, fg=TEXT_COLOR, relief="flat", padx=20, pady=10,
                     font=("Arial", 12, "bold"))
btn_file.pack(pady=10)
btn_file.bind("<Enter>", on_enter)
btn_file.bind("<Leave>", on_leave)

# Zone de texte avec cadre
text_frame = tk.Frame(main_frame, bg=BG_COLOR)
text_frame.pack(fill=tk.BOTH, expand=True, pady=10)

text_label = tk.Label(text_frame, text="Collez vos paroles ici :", 
                      font=("Arial", 11), fg=TEXT_COLOR, bg=BG_COLOR)
text_label.pack(anchor="w", pady=(0, 5))

text_box = tk.Text(text_frame, wrap="word", width=80, height=20, 
                   font=("Arial", 12), bg=ENTRY_BG, fg=TEXT_COLOR,
                   relief="solid", bd=1, padx=10, pady=10)
text_box.pack(fill=tk.BOTH, expand=True)

# Bouton de génération
btn_text = tk.Button(main_frame, text="✨ Générer le PowerPoint", command=generate_from_text,
                     bg=BUTTON_COLOR, fg=TEXT_COLOR, relief="flat", padx=20, pady=12,
                     font=("Arial", 12, "bold"))
btn_text.pack(pady=20)
btn_text.bind("<Enter>", on_enter)
btn_text.bind("<Leave>", on_leave)

root.mainloop()