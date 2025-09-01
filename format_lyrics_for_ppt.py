# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from datetime import date
import os

def split_block_recursive(block, max_lines=12):
    if len(block) <= max_lines:
        return [block]
    mid = len(block) // 2
    return split_block_recursive(block[:mid], max_lines) + split_block_recursive(block[mid:], max_lines)

# Lecture du fichier paroles.txt
with open("paroles.txt", "r", encoding="utf-8") as f:
    lines = f.readlines()

prs = Presentation()
block = []

for line in lines:
    line = line.strip()
    if line == "":
        if block:
            small_blocks = split_block_recursive(block, max_lines=8)
            for b in small_blocks:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title_placeholder = slide.shapes.title
                title_placeholder.text = "\n".join(b)

                text_frame = title_placeholder.text_frame
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                for p in text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    run = p.runs[0]
                    run.font.size = Pt(38)
                    run.font.bold = False
                    run.font.color.rgb = RGBColor(0, 0, 0)

                # Fond presque blanc avec légère nuance rose
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 240, 247)  # #fff0f7
            block = []
    else:
        block.append(line)

# Dernier bloc
if block:
    small_blocks = split_block_recursive(block, max_lines=8)
    for b in small_blocks:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_placeholder = slide.shapes.title
        title_placeholder.text = "\n".join(b)

        text_frame = title_placeholder.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = Pt(38)
            run.font.bold = False
            run.font.color.rgb = RGBColor(0, 0, 0)

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 240, 247)  # #fff0f7

# Nom du fichier avec date du jour au format DD-MM-YYYY
today = date.today().strftime("%d-%m-%Y")
filename = f"Chorale_{today}.pptx"
prs.save(filename)
print(f"PowerPoint créé : {filename}")
